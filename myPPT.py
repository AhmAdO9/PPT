import collections
import collections.abc
from pptx import Presentation
from pathlib import Path
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE




pr = Presentation()
slide_register1 = pr.slide_layouts[1]
slide1 = pr.slides.add_slide(slide_register1)
slide2 = pr.slides.add_slide(slide_register1)
text1_title = slide1.shapes.title
text1_content = slide1.placeholders[1]
text2_title = slide2.placeholders[0]
text2_content = slide2.placeholders[1]
path1 = Path('sample_slide1_input.txt').read_text()
path2 = Path('sample_slide2_input.txt').read_text()
text1_title.text= "SLIDE_ONE" 
text1_content.text = path1
for shape in slide1.shapes:  
    pass        

    textFrame = shape.text_frame                    
    textFrame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
    textFrame.word_wrap = True  
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(10) 
        paragraph.font.name = 'Love Ya Like A Sister'

    
 
text2_title.text= 'SLIDE_TWO'
text2_content.text = path2
for shape in slide2.shapes:  
    pass        

    textFrame = shape.text_frame                    
    textFrame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
    textFrame.word_wrap = True  
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(10) 
        paragraph.font.name = 'Love Ya Like A Sister'



pr.save("PPT.pptx")

